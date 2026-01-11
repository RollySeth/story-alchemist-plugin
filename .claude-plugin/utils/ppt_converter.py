"""
PowerPoint File Format Converter
Handles various PowerPoint formats (.ppt, .pptx, .pdf) and converts them to a readable format.
"""

import os
import sys
import tempfile
import zipfile
from pathlib import Path


def is_modern_pptx(file_path):
    """Check if file is a modern PPTX (ZIP-based) format."""
    try:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            return '[Content_Types].xml' in zip_ref.namelist()
    except (zipfile.BadZipFile, FileNotFoundError):
        return False


def is_legacy_ppt(file_path):
    """Check if file is a legacy PPT (OLE/Compound Document) format."""
    try:
        with open(file_path, 'rb') as f:
            header = f.read(8)
            # Check for OLE/Compound Document magic bytes
            return header[:8] == b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1'
    except FileNotFoundError:
        return False


def convert_legacy_ppt_to_pptx(ppt_path):
    """Convert legacy .ppt to modern .pptx using PowerPoint COM automation (Windows only)."""
    try:
        import win32com.client
        import pythoncom

        pythoncom.CoInitialize()

        try:
            ppt_app = win32com.client.Dispatch('PowerPoint.Application')

            # Create temp file for conversion
            temp_dir = tempfile.gettempdir()
            temp_pptx = os.path.join(temp_dir, f"converted_{os.path.basename(ppt_path)}.pptx")

            # Open the legacy PPT file
            presentation = ppt_app.Presentations.Open(
                os.path.abspath(ppt_path),
                ReadOnly=True,
                Untitled=True,
                WithWindow=False
            )

            # Save as modern PPTX format (24 = ppSaveAsOpenXMLPresentation)
            presentation.SaveAs(os.path.abspath(temp_pptx), 24)

            presentation.Close()
            ppt_app.Quit()

            return temp_pptx

        finally:
            pythoncom.CoUninitialize()

    except ImportError:
        print("ERROR: win32com not available. Install with: pip install pywin32", file=sys.stderr)
        return None
    except Exception as e:
        print(f"ERROR: Failed to convert legacy PPT: {e}", file=sys.stderr)
        return None


def extract_pptx_content(pptx_path):
    """Extract text and metadata from modern PPTX file."""
    try:
        from pptx import Presentation

        presentation = Presentation(pptx_path)

        result = {
            'total_slides': len(presentation.slides),
            'slides': []
        }

        for i, slide in enumerate(presentation.slides, 1):
            slide_data = {
                'slide_number': i,
                'shapes': []
            }

            for shape in slide.shapes:
                shape_info = {
                    'type': shape.shape_type
                }

                if hasattr(shape, 'text') and shape.text.strip():
                    shape_info['text'] = shape.text.strip()

                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    if shape.text_frame.text.strip():
                        shape_info['text'] = shape.text_frame.text.strip()

                        # Get font info if available
                        try:
                            if shape.text_frame.paragraphs:
                                first_run = shape.text_frame.paragraphs[0].runs[0] if shape.text_frame.paragraphs[0].runs else None
                                if first_run:
                                    shape_info['font'] = {
                                        'name': first_run.font.name,
                                        'size': first_run.font.size.pt if first_run.font.size else None
                                    }
                        except:
                            pass

                slide_data['shapes'].append(shape_info)

            result['slides'].append(slide_data)

        return result

    except ImportError:
        print("ERROR: python-pptx not available. Install with: pip install python-pptx", file=sys.stderr)
        return None
    except Exception as e:
        print(f"ERROR: Failed to extract PPTX content: {e}", file=sys.stderr)
        return None


def extract_ppt_content_via_com(ppt_path):
    """Extract content from PowerPoint file using COM automation (works for both .ppt and .pptx)."""
    try:
        import win32com.client
        import pythoncom

        pythoncom.CoInitialize()

        try:
            ppt_app = win32com.client.Dispatch('PowerPoint.Application')

            presentation = ppt_app.Presentations.Open(
                os.path.abspath(ppt_path),
                ReadOnly=True,
                Untitled=True,
                WithWindow=False
            )

            result = {
                'total_slides': presentation.Slides.Count,
                'slide_width': presentation.PageSetup.SlideWidth,
                'slide_height': presentation.PageSetup.SlideHeight,
                'slides': []
            }

            for i in range(1, presentation.Slides.Count + 1):
                slide = presentation.Slides(i)
                slide_data = {
                    'slide_number': i,
                    'shapes': []
                }

                # Get background info
                try:
                    if slide.Background.Fill.Visible:
                        fill_type = slide.Background.Fill.Type
                        if fill_type == 1:  # Solid color
                            try:
                                rgb = slide.Background.Fill.ForeColor.RGB
                                r = rgb & 255
                                g = (rgb >> 8) & 255
                                b = (rgb >> 16) & 255
                                slide_data['background_color'] = f'RGB({r}, {g}, {b})'
                            except:
                                pass
                except:
                    pass

                # Extract shapes
                for j in range(1, slide.Shapes.Count + 1):
                    shape = slide.Shapes(j)
                    shape_info = {
                        'type': shape.Type
                    }

                    # Extract text
                    if hasattr(shape, 'HasTextFrame') and shape.HasTextFrame:
                        if shape.TextFrame.HasText:
                            text = shape.TextFrame.TextRange.Text.strip()
                            if text:
                                shape_info['text'] = text

                                # Get font info
                                try:
                                    shape_info['font'] = {
                                        'name': shape.TextFrame.TextRange.Font.Name,
                                        'size': shape.TextFrame.TextRange.Font.Size
                                    }
                                except:
                                    pass

                    slide_data['shapes'].append(shape_info)

                result['slides'].append(slide_data)

            presentation.Close()
            ppt_app.Quit()

            return result

        finally:
            pythoncom.CoUninitialize()

    except ImportError:
        print("ERROR: win32com not available. Install with: pip install pywin32", file=sys.stderr)
        return None
    except Exception as e:
        print(f"ERROR: Failed to extract content via COM: {e}", file=sys.stderr)
        return None


def export_slides_as_images(ppt_path, output_dir=None, slide_numbers=None):
    """Export PowerPoint slides as PNG images."""
    if output_dir is None:
        output_dir = tempfile.mkdtemp()

    os.makedirs(output_dir, exist_ok=True)

    try:
        import win32com.client
        import pythoncom

        pythoncom.CoInitialize()

        try:
            ppt_app = win32com.client.Dispatch('PowerPoint.Application')

            presentation = ppt_app.Presentations.Open(
                os.path.abspath(ppt_path),
                ReadOnly=True,
                Untitled=True,
                WithWindow=False
            )

            exported_files = []

            # Determine which slides to export
            if slide_numbers is None:
                slide_numbers = range(1, min(presentation.Slides.Count + 1, 11))  # First 10 slides

            for slide_num in slide_numbers:
                if slide_num <= presentation.Slides.Count:
                    output_path = os.path.join(output_dir, f'slide_{slide_num}.png')
                    presentation.Slides(slide_num).Export(output_path, 'PNG', 1920, 1080)
                    exported_files.append(output_path)

            presentation.Close()
            ppt_app.Quit()

            return exported_files

        finally:
            pythoncom.CoUninitialize()

    except ImportError:
        print("ERROR: win32com not available. Install with: pip install pywin32", file=sys.stderr)
        return []
    except Exception as e:
        print(f"ERROR: Failed to export slides: {e}", file=sys.stderr)
        return []


def process_presentation(file_path, export_images=False, image_slides=None):
    """
    Main function to process any PowerPoint file format.

    Args:
        file_path: Path to the PowerPoint file (.ppt, .pptx, or .pdf)
        export_images: Whether to export slides as images
        image_slides: List of slide numbers to export (default: first 5)

    Returns:
        Dictionary with presentation content and metadata
    """
    file_path = os.path.abspath(file_path)

    if not os.path.exists(file_path):
        return {'error': f'File not found: {file_path}'}

    # Detect file format
    if is_modern_pptx(file_path):
        print(f"Detected modern PPTX format", file=sys.stderr)
        content = extract_pptx_content(file_path)
        if content is None:
            # Fallback to COM if python-pptx fails
            print("Falling back to COM automation...", file=sys.stderr)
            content = extract_ppt_content_via_com(file_path)

    elif is_legacy_ppt(file_path):
        print(f"Detected legacy PPT format", file=sys.stderr)
        # Use COM automation directly for legacy formats
        content = extract_ppt_content_via_com(file_path)

    else:
        # Try COM automation as last resort
        print(f"Unknown format, attempting COM automation...", file=sys.stderr)
        content = extract_ppt_content_via_com(file_path)

    if content is None:
        return {'error': 'Failed to extract presentation content'}

    # Export images if requested
    if export_images:
        if image_slides is None:
            image_slides = [1, 2, 3, 5, content.get('total_slides', 10)]

        exported_images = export_slides_as_images(file_path, slide_numbers=image_slides)
        content['exported_images'] = exported_images

    return content


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python ppt_converter.py <path_to_ppt_file> [--export-images]")
        sys.exit(1)

    file_path = sys.argv[1]
    export_images = '--export-images' in sys.argv

    result = process_presentation(file_path, export_images=export_images)

    if 'error' in result:
        print(f"ERROR: {result['error']}", file=sys.stderr)
        sys.exit(1)

    import json
    print(json.dumps(result, indent=2))
